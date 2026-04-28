"""Showcase 01 — Design system: tokens + recipes.

Builds a six-slide branded deck using the high-level recipes:

    title_slide → kpi_slide → bullet_slide → quote_slide
    → image_hero_slide → custom Grid-laid card row

Demonstrates: ``DesignTokens``, every recipe in ``pptx.design.recipes``,
the ``Grid`` layout primitive, and the lint-on-save safety net.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.design.layout import Grid
from pptx.design.recipes import (
    bullet_slide,
    image_hero_slide,
    kpi_slide,
    quote_slide,
    title_slide,
)
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from _lint import lint_or_die
from _tokens import BRAND

HERE = Path(__file__).parent
HERO = HERE / "assets" / "hero.jpg"


def build(out_path: Path) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide(
        prs,
        title="The power-pptx Showcase",
        subtitle="Space-aware authoring, design tokens & visual effects",
        tokens=BRAND,
    )

    kpi_slide(
        prs,
        title="Why this fork exists",
        kpis=[
            {"label": "Overflow bugs caught",  "value": "100%", "delta": +1.00},
            {"label": "Off-slide shapes auto-fixed", "value": "Yes", "delta": +0.95},
            {"label": "Manual cleanup needed", "value": "Zero",  "delta": -1.00},
        ],
        tokens=BRAND,
    )

    bullet_slide(
        prs,
        title="What ships in 1.1",
        bullets=[
            "Pre-flight text fitting via Pillow font metrics",
            "Slide linter with auto-fix for off-slide shapes",
            "Design tokens, slide recipes, Grid + Stack layout",
            "Animations, transitions, and motion paths",
            "Visual effects: shadows, glows, gradients, alpha",
        ],
        tokens=BRAND,
    )

    quote_slide(
        prs,
        quote="The new dashboards saved my team a week per sprint.",
        attribution="Director of Engineering, Flagship Customer",
        tokens=BRAND,
    )

    image_hero_slide(
        prs,
        title="Built for dynamic generation",
        image=str(HERO),
        tokens=BRAND,
    )

    # Custom Grid-laid card row to round out the deck.
    feature_slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    _add_section_title(feature_slide, "The pillars")
    grid = Grid(feature_slide, cols=12, rows=6,
                gutter=Pt(16), margin=Pt(48))
    features = [
        ("Space-aware", "Text never overflows. Shapes never wander off-slide."),
        ("Design tokens", "One palette + typography spec drives every recipe."),
        ("Visual polish", "Shadows, glows, gradients without writing OOXML."),
    ]
    for i, (heading, body) in enumerate(features):
        col = i * 4
        card = feature_slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 0, 0, 1, 1,
        )
        grid.place(card, col=col, row=2, col_span=4, row_span=3)
        _style_feature_card(card, heading, body)

    lint_or_die(prs)
    prs.save(out_path)
    return prs


def _add_section_title(slide, text: str) -> None:
    box = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.6), Inches(12), Inches(1.0),
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = text
    tf.fit_text(font_family="Inter", max_size=40, bold=True)
    tf.paragraphs[0].font.color.rgb = RGBColor(0x0F, 0x17, 0x2A)


def _style_feature_card(card, heading: str, body: str) -> None:
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    card.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    card.line.width = Pt(1)
    card.shadow.blur_radius = Pt(18)
    card.shadow.distance = Pt(4)
    card.shadow.direction = 90.0
    card.shadow.color.rgb = RGBColor(0x0F, 0x17, 0x2A)
    card.shadow.color.alpha = 0.10

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(18)
    tf.margin_top = tf.margin_bottom = Pt(18)

    tf.text = heading
    p0 = tf.paragraphs[0]
    p0.font.name = "Inter"
    p0.font.size = Pt(22)
    p0.font.bold = True
    p0.font.color.rgb = RGBColor(0x4F, 0x46, 0xE5)

    p1 = tf.add_paragraph()
    p1.text = body
    p1.font.name = "Inter"
    p1.font.size = Pt(14)
    p1.font.color.rgb = RGBColor(0x33, 0x41, 0x55)
    p1.space_before = Pt(8)


if __name__ == "__main__":
    out = HERE / "_out" / "01_design_system.pptx"
    out.parent.mkdir(exist_ok=True)
    build(out)
    print(f"wrote {out}")
