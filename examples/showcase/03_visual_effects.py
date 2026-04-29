"""Showcase 03 — Visual effects: shadows, glows, gradients, alpha.

Each slide isolates one effect family so the rendered thumbnails make
the difference legible at a glance.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.design.layout import Grid
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from _lint import lint_or_die
from _tokens import BRAND

HERE = Path(__file__).parent

PRIMARY = RGBColor(0x4F, 0x46, 0xE5)
ACCENT = RGBColor(0x22, 0xD3, 0xEE)
NEUTRAL = RGBColor(0x0F, 0x17, 0x2A)
SURFACE = RGBColor(0xF8, 0xFA, 0xFC)


def build(out_path: Path) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _shadow_slide(prs)
    _glow_slide(prs)
    _gradient_slide(prs)
    _alpha_slide(prs)

    lint_or_die(prs)
    prs.save(out_path)
    return prs


def _add_title(slide, text: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = text
    tf.fit_text(font_family="Inter", max_size=30, bold=True)
    tf.paragraphs[0].font.color.rgb = NEUTRAL

    if subtitle:
        sb = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.15), Inches(12), Inches(0.5),
        )
        sb.text_frame.text = subtitle
        sb.text_frame.paragraphs[0].font.size = Pt(14)
        sb.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x64, 0x74, 0x8B)


def _grid_card(slide, grid: Grid, *, col: int, row: int,
               col_span: int = 4, row_span: int = 3):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 0, 0, 1, 1)
    grid.place(card, col=col, row=row, col_span=col_span, row_span=row_span)
    card.fill.solid()
    card.fill.fore_color.rgb = SURFACE
    card.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    card.line.width = Pt(1)
    return card


def _label(card, heading: str, body: str) -> None:
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(18)
    tf.margin_top = tf.margin_bottom = Pt(18)
    tf.text = heading
    p0 = tf.paragraphs[0]
    p0.font.name = "Inter"
    p0.font.size = Pt(20)
    p0.font.bold = True
    p0.font.color.rgb = NEUTRAL
    p1 = tf.add_paragraph()
    p1.text = body
    p1.font.name = "Inter"
    p1.font.size = Pt(13)
    p1.font.color.rgb = RGBColor(0x47, 0x55, 0x69)
    p1.space_before = Pt(6)


def _shadow_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Outer shadow", "Three intensities of the same drop shadow.")
    grid = Grid(slide, cols=12, rows=6, gutter=Pt(20), margin=Pt(48))

    presets = [
        ("Subtle", "blur 8 · α 12%", 8.0, 2.0, 0.12),
        ("Card",   "blur 18 · α 18%", 18.0, 4.0, 0.18),
        ("Hero",   "blur 36 · α 28%", 36.0, 8.0, 0.28),
    ]
    for i, (heading, body, blur, dist, alpha) in enumerate(presets):
        card = _grid_card(slide, grid, col=i * 4, row=2)
        card.shadow.blur_radius = Pt(blur)
        card.shadow.distance = Pt(dist)
        card.shadow.direction = 90.0
        card.shadow.color.rgb = NEUTRAL
        card.shadow.color.alpha = alpha
        _label(card, heading, body)


def _glow_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Glow & soft edges",
               "Useful for badges, callouts, and active states.")
    grid = Grid(slide, cols=12, rows=6, gutter=Pt(20), margin=Pt(48))

    cards = [
        ("Primary glow", "radius 8 · indigo",   PRIMARY, 8),
        ("Accent glow",  "radius 12 · cyan",    ACCENT,  12),
        ("Soft edges",   "soft-edge radius 6",  None,    None),
    ]
    for i, (heading, body, color, radius) in enumerate(cards):
        card = _grid_card(slide, grid, col=i * 4, row=2)
        if color is not None:
            card.glow.radius = Pt(radius)
            card.glow.color.rgb = color
        else:
            card.soft_edges.radius = Pt(6)
        _label(card, heading, body)


def _gradient_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Gradient fills",
               "Linear, radial, and shape gradients via mutable stops.")
    grid = Grid(slide, cols=12, rows=6, gutter=Pt(20), margin=Pt(48))

    # Linear (diagonal-ish — direction is rendered by PowerPoint)
    c1 = _grid_card(slide, grid, col=0, row=2)
    c1.fill.gradient(kind="linear")
    c1.fill.gradient_stops.replace([(0.0, "#4F46E5"), (1.0, "#22D3EE")])
    _label(c1, "Linear", "indigo → cyan")
    _force_white_text(c1)

    # Radial
    c2 = _grid_card(slide, grid, col=4, row=2)
    c2.fill.gradient(kind="radial")
    c2.fill.gradient_stops.replace([(0.0, "#22D3EE"), (1.0, "#0F172A")])
    _label(c2, "Radial", "cyan → slate")
    _force_white_text(c2)

    # Shape (matches the rounded-rectangle outline)
    c3 = _grid_card(slide, grid, col=8, row=2)
    c3.fill.gradient(kind="shape")
    c3.fill.gradient_stops.replace([(0.0, "#22D3EE"), (1.0, "#4F46E5")])
    _label(c3, "Shape", "follows the rectangle path")
    _force_white_text(c3)


def _alpha_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Alpha-tinted fills",
               "Glassy cards over a brand-colored backdrop.")

    # Backdrop band
    backdrop = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(1.7), prs.slide_width, Inches(5.0),
    )
    backdrop.fill.solid()
    backdrop.fill.fore_color.rgb = PRIMARY
    backdrop.line.fill.background()

    grid = Grid(slide, cols=12, rows=6, gutter=Pt(20), margin=Pt(48))

    presets = [
        ("Glass 25%", "indigo · α 0.25",  0.25),
        ("Glass 50%", "indigo · α 0.50",  0.50),
        ("Glass 75%", "indigo · α 0.75",  0.75),
    ]
    for i, (heading, body, alpha) in enumerate(presets):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 0, 0, 1, 1)
        grid.place(card, col=i * 4, row=2, col_span=4, row_span=3)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        card.fill.fore_color.alpha = alpha
        card.line.fill.background()
        card.shadow.blur_radius = Pt(24)
        card.shadow.distance = Pt(6)
        card.shadow.color.rgb = NEUTRAL
        card.shadow.color.alpha = 0.20
        _label(card, heading, body)


def _force_white_text(card) -> None:
    tf = card.text_frame
    for p in tf.paragraphs:
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


if __name__ == "__main__":
    out = HERE / "_out" / "03_visual_effects.pptx"
    out.parent.mkdir(exist_ok=True)
    build(out)
    print(f"wrote {out}")
