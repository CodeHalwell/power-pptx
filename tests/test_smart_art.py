"""Unit-test suite for `pptx.smart_art` module."""

from __future__ import annotations

import io
import os

import pytest

FIXTURE_PATH = os.path.join(os.path.dirname(__file__), "test_files", "smart_art_org_chart.pptx")


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _make_prs():
    """Return a fresh Presentation loaded from the SmartArt fixture."""
    from pptx import Presentation

    return Presentation(FIXTURE_PATH)


# ---------------------------------------------------------------------------
# SmartArtCollection
# ---------------------------------------------------------------------------


class DescribeSmartArtCollection:
    def it_has_length_1_for_the_org_chart_fixture(self):
        prs = _make_prs()
        assert len(prs.slides[0].smart_art) == 1

    def it_is_empty_on_a_blank_slide(self):
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        assert len(slide.smart_art) == 0

    def it_supports_iteration(self):
        prs = _make_prs()
        items = list(prs.slides[0].smart_art)
        assert len(items) == 1

    def it_raises_on_out_of_range_index(self):
        prs = _make_prs()
        with pytest.raises(IndexError):
            _ = prs.slides[0].smart_art[99]

    def it_returns_the_correct_repr(self):
        prs = _make_prs()
        r = repr(prs.slides[0].smart_art)
        assert "SmartArtCollection" in r
        assert "1" in r


# ---------------------------------------------------------------------------
# SmartArtShape.texts
# ---------------------------------------------------------------------------


class DescribeSmartArtShape_texts:
    def it_returns_the_initial_node_texts(self):
        prs = _make_prs()
        sa = prs.slides[0].smart_art[0]
        assert sa.texts == ["CEO", "CTO", "CFO"]

    def it_has_the_correct_shape_name(self):
        prs = _make_prs()
        sa = prs.slides[0].smart_art[0]
        assert sa.name == "SmartArt 1"


# ---------------------------------------------------------------------------
# SmartArtShape.set_text
# ---------------------------------------------------------------------------


class DescribeSmartArtShape_set_text:
    def it_replaces_all_node_texts(self):
        prs = _make_prs()
        sa = prs.slides[0].smart_art[0]
        sa.set_text(["Alice", "Bob", "Carol"])
        assert sa.texts == ["Alice", "Bob", "Carol"]

    def it_round_trips_through_save_and_reopen(self):
        from pptx import Presentation

        prs = _make_prs()
        prs.slides[0].smart_art[0].set_text(["Alice", "Bob", "Carol"])

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        prs2 = Presentation(buf)

        assert prs2.slides[0].smart_art[0].texts == ["Alice", "Bob", "Carol"]

    def it_raises_on_wrong_length_in_strict_mode(self):
        prs = _make_prs()
        sa = prs.slides[0].smart_art[0]
        with pytest.raises(ValueError, match="3 content node"):
            sa.set_text(["Only", "Two"])

    def it_silently_accepts_extra_values_when_strict_is_False(self):
        prs = _make_prs()
        sa = prs.slides[0].smart_art[0]
        # 5 values for 3 nodes — excess should be ignored
        sa.set_text(["A", "B", "C", "D", "E"], strict=False)
        assert sa.texts == ["A", "B", "C"]

    def it_accepts_fewer_values_when_strict_is_False(self):
        prs = _make_prs()
        sa = prs.slides[0].smart_art[0]
        sa.set_text(["X"], strict=False)
        # only first node changed; rest remain
        texts = sa.texts
        assert texts[0] == "X"
        assert texts[1] == "CTO"
        assert texts[2] == "CFO"

    def it_leaves_layout_colors_and_style_parts_unchanged(self):
        """The three non-data parts should be byte-for-byte identical after set_text."""
        from pptx import Presentation

        prs = _make_prs()

        # Capture original sibling part bytes
        slide_part = prs.slides[0].part
        lo_part = slide_part.related_part("rId2")
        qs_part = slide_part.related_part("rId3")
        cs_part = slide_part.related_part("rId4")
        orig_lo = lo_part.blob
        orig_qs = qs_part.blob
        orig_cs = cs_part.blob

        # Mutate the data part
        prs.slides[0].smart_art[0].set_text(["X", "Y", "Z"])

        # Re-fetch the parts — should be unchanged
        assert slide_part.related_part("rId2").blob == orig_lo
        assert slide_part.related_part("rId3").blob == orig_qs
        assert slide_part.related_part("rId4").blob == orig_cs
