"""Unit tests for `pptx.design.layout`."""

from __future__ import annotations

import pytest

from pptx import Presentation
from pptx.design.layout import Box, Grid, Stack
from pptx.util import Emu, Inches, Pt


@pytest.fixture
def slide():
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])


class DescribeBox:
    def it_unpacks_to_left_top_width_height(self):
        box = Box(Emu(1), Emu(2), Emu(3), Emu(4))
        left, top, width, height = box
        assert (left, top, width, height) == (1, 2, 3, 4)


class DescribeGrid:
    def it_exposes_cols_and_rows(self, slide):
        grid = Grid(slide, cols=12, rows=6)

        assert grid.cols == 12
        assert grid.rows == 6

    def it_returns_a_box_for_a_single_cell(self, slide):
        grid = Grid(slide, cols=12, rows=6, gutter=Pt(12), margin=Inches(0.5))

        box = grid.cell(col=0, row=0)

        assert isinstance(box, Box)
        assert box.left == Inches(0.5)
        assert box.top == Inches(0.5)
        # cell width = (slide_w - 2*margin - 11*gutter) / 12
        slide_w = 9144000
        usable_w = slide_w - 2 * int(Inches(0.5)) - 11 * int(Pt(12))
        assert box.width == int(round(usable_w / 12))

    def it_spans_multiple_columns(self, slide):
        grid = Grid(slide, cols=12, rows=6, gutter=Pt(12), margin=Inches(0.5))

        full = grid.cell(col=0, row=0, col_span=12, row_span=1)
        half_left = grid.cell(col=0, row=0, col_span=6, row_span=1)
        half_right = grid.cell(col=6, row=0, col_span=6, row_span=1)

        # full span = sum of left half + gutter + right half
        assert full.width == half_left.width + int(Pt(12)) + half_right.width
        # right half starts where left half ends + 1 gutter
        assert half_right.left == half_left.left + half_left.width + int(Pt(12))

    def it_raises_on_out_of_bounds_cells(self, slide):
        grid = Grid(slide, cols=4, rows=2)
        with pytest.raises(IndexError):
            grid.cell(col=4, row=0)
        with pytest.raises(IndexError):
            grid.cell(col=0, row=2)
        with pytest.raises(IndexError):
            grid.cell(col=3, row=0, col_span=2)

    def it_rejects_negative_indices_and_zero_spans(self, slide):
        grid = Grid(slide, cols=4, rows=2)
        with pytest.raises(IndexError):
            grid.cell(col=-1, row=0)
        with pytest.raises(IndexError):
            grid.cell(col=0, row=0, col_span=0)

    def it_rejects_invalid_construction(self, slide):
        with pytest.raises(ValueError):
            Grid(slide, cols=0)
        with pytest.raises(ValueError):
            Grid(slide, cols=4, rows=0)

    def it_rejects_margins_that_consume_the_slide(self, slide):
        with pytest.raises(ValueError):
            Grid(slide, cols=4, margin=Inches(20))

    def it_raises_a_clear_error_for_a_detached_slide(self):
        # An object that looks slide-shaped but isn't attached to a presentation
        class _FakeSlide:
            pass

        with pytest.raises(ValueError, match="attached to a presentation"):
            Grid(_FakeSlide(), cols=4)

    def it_accepts_2_tuple_and_4_tuple_margins(self, slide):
        g2 = Grid(slide, cols=2, rows=1, margin=(Pt(10), Pt(20)))
        # 2-tuple is (vertical, horizontal)
        assert g2.cell(col=0, row=0).left == int(Pt(20))
        assert g2.cell(col=0, row=0).top == int(Pt(10))

        g4 = Grid(slide, cols=2, rows=1, margin=(Pt(1), Pt(2), Pt(3), Pt(4)))
        # 4-tuple is (top, right, bottom, left)
        assert g4.cell(col=0, row=0).left == int(Pt(4))
        assert g4.cell(col=0, row=0).top == int(Pt(1))

    def it_can_place_a_shape_into_a_cell(self, slide):
        grid = Grid(slide, cols=12, rows=6, gutter=Pt(12), margin=Inches(0.5))
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(1), Inches(1))

        returned = grid.place(shape, col=0, row=0, col_span=12, row_span=1)

        assert returned is shape
        expected = grid.cell(col=0, row=0, col_span=12, row_span=1)
        assert (shape.left, shape.top, shape.width, shape.height) == tuple(expected)


class DescribeStack:
    def it_lays_cells_out_top_to_bottom_for_vertical(self):
        stack = Stack(
            direction="vertical",
            gap=Pt(8),
            left=Inches(1),
            top=Inches(1),
            width=Inches(8),
        )

        b1 = stack.next(height=Inches(1))
        b2 = stack.next(height=Inches(2))

        assert b1 == Box(Inches(1), Inches(1), Inches(8), Inches(1))
        assert b2.top == b1.top + b1.height + int(Pt(8))
        assert b2.height == Inches(2)

    def it_lays_cells_out_left_to_right_for_horizontal(self):
        stack = Stack(
            direction="horizontal",
            gap=Pt(4),
            left=Inches(0),
            top=Inches(0),
            height=Inches(2),
        )

        b1 = stack.next(width=Inches(1))
        b2 = stack.next(width=Inches(2))

        assert b1.left == Inches(0)
        assert b2.left == b1.left + b1.width + int(Pt(4))
        assert b2.height == Inches(2)

    def it_does_not_insert_a_leading_gap_before_the_first_cell(self):
        stack = Stack(
            direction="vertical",
            gap=Pt(50),
            left=Inches(0),
            top=Inches(2),
            width=Inches(1),
        )
        first = stack.next(height=Inches(1))
        assert first.top == Inches(2)

    def it_supports_per_cell_width_override_for_vertical(self):
        stack = Stack(direction="vertical", left=Inches(0), top=Inches(0), width=Inches(8))
        box = stack.next(width=Inches(4), height=Inches(1))
        assert box.width == Inches(4)

    def it_resets_the_cursor(self):
        stack = Stack(direction="vertical", left=Inches(0), top=Inches(0), width=Inches(1))
        stack.next(height=Inches(1))
        stack.reset()
        b = stack.next(height=Inches(1))
        assert b.top == Inches(0)

    def it_rejects_invalid_directions(self):
        with pytest.raises(ValueError):
            Stack(direction="diagonal")

    def it_requires_height_for_vertical_next(self):
        stack = Stack(direction="vertical", width=Inches(1))
        with pytest.raises(TypeError):
            stack.next()

    def it_requires_width_for_horizontal_next(self):
        stack = Stack(direction="horizontal", height=Inches(1))
        with pytest.raises(TypeError):
            stack.next()

    def it_requires_a_width_for_vertical_stack(self):
        stack = Stack(direction="vertical")
        with pytest.raises(TypeError):
            stack.next(height=Inches(1))

    def it_requires_a_height_for_horizontal_stack(self):
        stack = Stack(direction="horizontal")
        with pytest.raises(TypeError):
            stack.next(width=Inches(1))

    def it_can_place_a_shape(self, slide):
        stack = Stack(
            direction="vertical",
            gap=Pt(8),
            left=Inches(0),
            top=Inches(0),
            width=Inches(2),
        )
        shape = slide.shapes.add_shape(1, Inches(5), Inches(5), Inches(1), Inches(1))

        returned = stack.place(shape, height=Inches(1))

        assert returned is shape
        assert (shape.left, shape.top) == (Inches(0), Inches(0))
        assert (shape.width, shape.height) == (Inches(2), Inches(1))
