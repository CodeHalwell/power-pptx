"""Unit-test suite for :mod:`pptx.design.style`."""

from __future__ import annotations

import pytest

from pptx import Presentation
from pptx.design.style import ShapeStyle
from pptx.design.tokens import ShadowToken, TypographyToken
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


@pytest.fixture
def shape():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(2), Inches(1))


class DescribeShapeStyle:
    def it_is_exposed_on_base_shape(self, shape):
        assert isinstance(shape.style, ShapeStyle)

    def it_returns_the_same_facade_each_read(self, shape):
        assert shape.style is shape.style

    def it_applies_a_solid_fill_from_an_rgb_color(self, shape):
        shape.style.fill = RGBColor(0xFF, 0x00, 0x00)
        assert shape.fill.fore_color.rgb == RGBColor(0xFF, 0x00, 0x00)

    def it_coerces_a_hex_string_fill(self, shape):
        shape.style.fill = "#3C2F80"
        assert shape.fill.fore_color.rgb == RGBColor(0x3C, 0x2F, 0x80)

    def it_clears_the_fill_when_assigned_none(self, shape):
        shape.style.fill = "#FF0000"
        shape.style.fill = None
        from pptx.enum.dml import MSO_FILL

        assert shape.fill.type == MSO_FILL.BACKGROUND

    def it_applies_a_line_color(self, shape):
        shape.style.line = (0x12, 0x34, 0x56)
        assert shape.line.color.rgb == RGBColor(0x12, 0x34, 0x56)

    def it_applies_a_shadow_token(self, shape):
        token = ShadowToken(
            blur_radius=Pt(8),
            distance=Pt(2),
            direction=90.0,
            color=RGBColor(0, 0, 0),
            alpha=0.25,
        )
        shape.style.shadow = token
        assert shape.shadow.blur_radius == Pt(8)
        assert shape.shadow.distance == Pt(2)
        assert shape.shadow.direction == 90.0

    def it_leaves_unset_shadow_fields_alone(self, shape):
        shape.shadow.distance = Pt(5)
        shape.style.shadow = ShadowToken(blur_radius=Pt(8))
        # blur was applied, distance preserved.
        assert shape.shadow.blur_radius == Pt(8)
        assert shape.shadow.distance == Pt(5)

    def it_clears_the_shadow_when_assigned_none(self, shape):
        shape.style.shadow = ShadowToken(blur_radius=Pt(8), distance=Pt(2))
        shape.style.shadow = None
        assert shape.shadow.blur_radius is None
        assert shape.shadow.distance is None

    def it_applies_a_font_token_to_every_run(self, shape):
        shape.text_frame.text = "hello"
        token = TypographyToken(family="Inter", size=Pt(20), bold=True)
        shape.style.font = token
        run = shape.text_frame.paragraphs[0].runs[0]
        assert run.font.name == "Inter"
        assert run.font.size == Pt(20)
        assert run.font.bold is True

    def it_applies_text_color_to_every_run(self, shape):
        shape.text_frame.text = "hello"
        shape.style.text_color = "#3C2F80"
        run = shape.text_frame.paragraphs[0].runs[0]
        assert run.font.color.rgb == RGBColor(0x3C, 0x2F, 0x80)

    def it_rejects_reading_writeonly_attrs(self, shape):
        with pytest.raises(AttributeError):
            shape.style.fill
        with pytest.raises(AttributeError):
            shape.style.line
        with pytest.raises(AttributeError):
            shape.style.shadow
